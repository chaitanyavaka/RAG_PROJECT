from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_ppt():
    prs = Presentation()

    # Slide 1: Title
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Agentic RAG Chatbot"
    subtitle.text = "Architecture & Implementation\nUsing Model Context Protocol (MCP)"

    # Slide 2: Architecture
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Agent-Based Architecture"
    content = slide.placeholders[1]
    content.text = (
        "1. Coordinator Agent: Orchestrates the workflow.\n"
        "2. Ingestion Agent: Parses PDF, DOCX, CSV, PPTX.\n"
        "3. Retrieval Agent: Semantic search with ChromaDB (all-MiniLM-L6-v2).\n"
        "4. LLM Response Agent: Generates answers using Groq (Llama-3).\n\n"
        "Communication via MCP (Model Context Protocol)."
    )

    # Slide 3: System Flow
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "System Flow"
    content = slide.placeholders[1]
    content.text = (
        "User Query -> Coordinator -> Retrieval Agent -> Context Found\n"
        "Coordinator -> LLM Agent (Context + Query) -> Answer\n"
        "Coordinator -> User UI"
    )

    # Slide 4: Tech Stack
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Technology Stack"
    content = slide.placeholders[1]
    content.text = (
        "- Backend: Python, FastAPI\n"
        "- Frontend: HTML5, TailwindCSS, Vanilla JS (Interactive)\n"
        "- Database: ChromaDB (Vector Store)\n"
        "- AI/LLM: Groq API (Llama-3-70b)\n"
        "- Protocol: MCP (Custom Implementation)"
    )

    # Slide 5: Challenges & Future
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Challenges & Improvements"
    content = slide.placeholders[1]
    content.text = (
        "Challenges:\n"
        "- Async agent coordination.\n"
        "- Parsing complex PDF layouts.\n\n"
        "Future Scope:\n"
        "- Distributed agents using RabbitMQ.\n"
        "- Image/Multi-modal support.\n"
        "- Persistent MCP connection over WebSocket."
    )

    prs.save('Agentic_RAG_Architecture.pptx')
    print("PPTX saved as Agentic_RAG_Architecture.pptx")

if __name__ == "__main__":
    create_ppt()
