import os
import io
import pandas as pd
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from .base import BaseAgent
from ..mcp.protocol import MCPMessage, MessageType

class IngestionAgent(BaseAgent):
    def __init__(self):
        super().__init__("IngestionAgent")
    
    async def process_message(self, message: MCPMessage):
        if message.type == MessageType.TASK_REQUEST:
            task = message.payload.get("task")
            if task == "ingest_file":
                await self.ingest_file(message)

    async def ingest_file(self, message: MCPMessage):
        file_path = message.payload.get("file_path")
        file_name = message.payload.get("file_name")
        
        try:
            text = self.extract_text(file_path, file_name)
            chunks = self.chunk_text(text)
            
            # Send chunks to RetrievalAgent for embedding
            await self.send_message(
                receiver="RetrievalAgent",
                type=MessageType.TASK_REQUEST,
                payload={
                    "task": "embed_chunks",
                    "chunks": chunks,
                    "metadata": {"source": file_name}
                },
                trace_id=message.trace_id
            )
            
            # Notify Coordinator/User of success
            await self.send_message(
                receiver=message.sender,
                type=MessageType.TASK_RESULT,
                payload={"status": "success", "file": file_name, "chunks_count": len(chunks)},
                trace_id=message.trace_id
            )
        except Exception as e:
            await self.send_message(
                receiver=message.sender,
                type=MessageType.ERROR,
                payload={"error": str(e)},
                trace_id=message.trace_id
            )

    def extract_text(self, file_path: str, file_name: str) -> str:
        ext = os.path.splitext(file_name)[1].lower()
        text = ""
        
        if ext == ".pdf":
            reader = PdfReader(file_path)
            for page in reader.pages:
                text += page.extract_text() + "\n"
        elif ext == ".docx":
            doc = Document(file_path)
            for para in doc.paragraphs:
                text += para.text + "\n"
        elif ext == ".txt" or ext == ".md":
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read()
        elif ext == ".csv":
            df = pd.read_csv(file_path)
            text = df.to_string()
        elif ext == ".pptx":
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        else:
            raise ValueError(f"Unsupported file format: {ext}")
            
        return text

    def chunk_text(self, text: str, chunk_size: int = 500, overlap: int = 50) -> list[str]:
        # Simple character-based chunking for now
        chunks = []
        for i in range(0, len(text), chunk_size - overlap):
            chunks.append(text[i:i + chunk_size])
        return chunks
